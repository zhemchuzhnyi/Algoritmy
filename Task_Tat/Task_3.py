from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Создаем презентацию
prs = Presentation()

# Слайд 1: Титульный
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Анализ заявок на сервис Докма"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
subtitle.text = "Период: 12.05.2025 – 25.05.2025\nПодготовлено: Grok\nДата: 28 мая 2025"
subtitle.text_frame.paragraphs[0].font.size = Pt(20)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)

# Слайд 2: Обзор данных
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Обзор данных"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
content = slide.placeholders[1].text_frame
content.text = "Умеренный режим (12.05–18.05):\n"
content.paragraphs[0].font.size = Pt(18)
content.add_paragraph().text = "• Общее количество заявок: 731"
content.add_paragraph().text = "• Среднее в день: ~104"
content.add_paragraph().text = "• Пики: 15.05 (134), 13.05 (121)"
content.add_paragraph().text = "• Спады: 17.05 (72), 18.05 (67)"
content.add_paragraph().text = "\nРедкий режим (19.05–25.05):\n"
content.add_paragraph().text = "• Общее количество заявок: 525"
content.add_paragraph().text = "• Среднее в день: ~75"
content.add_paragraph().text = "• Пики: 23.05 (86), 19.05 (85)"
content.add_paragraph().text = "• Спады: 24.05 (56), 25.05 (68)"
for p in content.paragraphs[1:]:
    p.font.size = Pt(18)
    p.level = 1 if p.text.startswith("•") else 0

# Слайд 3: График
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "График количества заявок по дням"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
img_path = 'applications_chart.png'
slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9))

# Слайд 4: Выводы
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Выводы"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
content = slide.placeholders[1].text_frame
content.text = "• Снижение активности: во второй неделе заявок на ~28% меньше (525 против 731)."
content.add_paragraph().text = "• Пики в середине недели: 13.05 (121), 15.05 (134) в умеренном режиме; 19.05 (85), 23.05 (86) в редком — вероятно, рабочие дни."
content.add_paragraph().text = "• Спад в выходные: 17–18.05 (72, 67) и 24–25.05 (56, 68) — меньшая активность пользователей."
for p in content.paragraphs:
    p.font.size = Pt(18)
    p.level = 1 if p.text.startswith("•") else 0

# Сохраняем презентацию
prs.save('applications_presentation.pptx')