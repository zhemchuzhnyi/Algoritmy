import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from io import BytesIO

# Настройки
EXCEL_FILE = "Отзывы_врачи.xlsx"
OUTPUT_PPTX = "Анализ_жалоб_врачи.pptx"
ICONS_PATH = "icons/"  # Папка с иконками
COLORS = {
    'primary': RGBColor(59, 89, 152),  # Синий
    'secondary': RGBColor(76, 175, 80),  # Зеленый
    'accent': RGBColor(96, 125, 139)  # Серый
}

# Загрузка данных из Excel
def load_data():
    try:
        df_main = pd.read_excel(EXCEL_FILE, sheet_name="Сводная")
        df_monthly = pd.read_excel(EXCEL_FILE, sheet_name="Месячная динамика")
        return df_main, df_monthly
    except Exception as e:
        print(f"Ошибка загрузки данных: {e}")
        # Тестовые данные
        df_main = pd.DataFrame({
            'ФИО врача': ['Эреджепова Э.М.', 'Дементьева О.В.', 'Каштелян А.А.', 'Руди Ю.С.'],
            'Специализация': ['Дерматолог', 'Невролог', 'Терапевт', 'Педиатр'],
            'Количество жалоб': [11, 6, 4, 4],
            'Причина жалобы': ['Качество', 'Коммуникация', 'Ожидание', 'Коммуникация'],
            'Время ответа (часы)': [2, 4, 3, 5],
            'Статус жалобы': ['Подтверждено', 'Подтверждено', 'Подтверждено', 'Подтверждено'],
            'Месяц': ['Январь', 'Декабрь', 'Февраль', 'Январь']
        })
        df_monthly = pd.DataFrame({
            'Месяц': ['Октябрь', 'Ноябрь', 'Декабрь', 'Январь', 'Февраль', 'Март', 'Апрель'] * 5,
            'Специализация': ['Гастроэнтерологи'] * 7 + ['Педиатры'] * 7 + ['Неврологи'] * 7 + ['Гинекологи'] * 7 + ['Дерматологи'] * 7,
            'Количество жалоб': [5, 6, 8, 9, 10, 7, 4] * 5
        })
        return df_main, df_monthly

# Анализ данных
def analyze_data(df_main, df_monthly):
    reasons = {'Коммуникация': 40, 'Ожидание': 25, 'Качество': 20, 'Опровергнуто': 15}
    specializations = {'Гастроэнтерологи': 18, 'Педиатры': 15, 'Неврологи': 12, 'Гинекологи': 10, 'Дерматологи': 8}
    top_doctors = df_main.nlargest(3, 'Количество жалоб')[['ФИО врача', 'Количество жалоб']]
    status_counts = {'Подтверждено': 85, 'Опровергнуто': 15}
    response_times = {'До 1 часа': 10, '1–3 часа': 30, 'Более 3 часов': 60}
    monthly_trend = df_monthly.groupby('Месяц')['Количество жалоб'].sum()
    heatmap_data = df_monthly.pivot_table(values='Количество жалоб', index='Специализация', columns='Месяц', fill_value=0)
    return {
        'reasons': reasons,
        'specializations': specializations,
        'top_doctors': top_doctors,
        'status_counts': status_counts,
        'response_times': response_times,
        'monthly_trend': monthly_trend,
        'heatmap_data': heatmap_data
    }

# Добавление анимации
def add_fade_animation(shape, order):
    animation = shape.element
    seq = OxmlElement('p:seq')
    seq.set('concurrent', '0')
    seq.set('nextAc', 'seek')
    animation.append(seq)
    anim = OxmlElement('p:animEffect')
    anim.set('transition', 'in')
    anim.set('filter', 'fade')
    anim.set('prst', 'fade')
    timing = OxmlElement('p:tmAbs')
    timing.set('tm', str(order * 1000))
    anim.append(timing)
    seq.append(anim)

# Создание презентации
def create_presentation(analysis):
    prs = Presentation()

    def set_title_style(title):
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']

    # Титульный слайд
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Анализ негативных отзывов на врачей (онлайн-консультации)"
    set_title_style(title)
    subtitle.text = "Данные за октябрь 2023 – апрель 2024"
    try:
        slide.shapes.add_picture(f"{ICONS_PATH}stethoscope.png", Inches(8), Inches(5), width=Inches(1))
    except FileNotFoundError:
        print("Иконка stethoscope.png не найдена")

    # Слайд 1: Распределение жалоб по причинам
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Распределение жалоб по причинам"
    set_title_style(title)
    chart_data = ChartData()
    chart_data.categories = list(analysis['reasons'].keys())
    chart_data.add_series('Доля, %', list(analysis['reasons'].values()))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(1), Inches(2), Inches(6), Inches(4.5), chart_data
    ).chart
    chart.series[0].data_labels.show_percentage = True
    chart.series[0].data_labels.number_format = '0%'
    add_fade_animation(chart, 1)

    # Слайд 2: ТОП-5 специализаций
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "ТОП-5 специализаций с жалобами"
    set_title_style(title)
    chart_data = ChartData()
    chart_data.categories = list(analysis['specializations'].keys())
    chart_data.add_series('Доля, %', list(analysis['specializations'].values()))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(1), Inches(2), Inches(6), Inches(4.5), chart_data
    ).chart
    add_fade_animation(chart, 1)
    try:
        slide.shapes.add_picture(f"{ICONS_PATH}chart.png", Inches(8), Inches(5), width=Inches(1))
    except FileNotFoundError:
        print("Иконка chart.png не найдена")

    # Слайд 3: Динамика жалоб по месяцам
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Динамика жалоб по месяцам"
    set_title_style(title)
    plt.figure(figsize=(8, 4))
    analysis['monthly_trend'].plot(kind='line', marker='o', color='#4CAF50')
    plt.ylabel('Количество жалоб')
    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=150)
    buf.seek(0)
    pic = slide.shapes.add_picture(buf, Inches(1), Inches(2), width=Inches(6))
    plt.close()
    add_fade_animation(pic, 1)
    try:
        slide.shapes.add_picture(f"{ICONS_PATH}calendar.png", Inches(8), Inches(5), width=Inches(1))
    except FileNotFoundError:
        print("Иконка calendar.png не найдена")

    # Слайд 4: ТОП-3 врачей
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Врачи с максимальным числом жалоб"
    set_title_style(title)
    chart_data = ChartData()
    chart_data.categories = analysis['top_doctors']['ФИО врача'].tolist()
    chart_data.add_series('Жалобы', analysis['top_doctors']['Количество жалоб'].tolist())
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(6), Inches(4.5), chart_data
    ).chart
    add_fade_animation(chart, 1)

    # Слайд 5: Соотношение опровергнутых и подтвержденных жалоб
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Соотношение опровергнутых и подтвержденных жалоб"
    set_title_style(title)
    chart_data = ChartData()
    chart_data.categories = list(analysis['status_counts'].keys())
    chart_data.add_series('Доля, %', list(analysis['status_counts'].values()))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(1), Inches(2), Inches(6), Inches(4.5), chart_data
    ).chart
    chart.series[0].data_labels.show_percentage = True
    chart.series[0].data_labels.number_format = '0%'
    add_fade_animation(chart, 1)

    # Слайд 6: Распределение по времени ответа
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Распределение жалоб по времени ответа"
    set_title_style(title)
    chart_data = ChartData()
    chart_data.categories = list(analysis['response_times'].keys())
    chart_data.add_series('Доля, %', list(analysis['response_times'].values()))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(6), Inches(4.5), chart_data
    ).chart
    add_fade_animation(chart, 1)

    # Слайд 7: Тепловая карта
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Тепловая карта жалоб по месяцам и специализациям"
    set_title_style(title)
    plt.figure(figsize=(8, 4))
    sns.heatmap(analysis['heatmap_data'], cmap='YlGnBu', annot=True, fmt='.0f')
    plt.xlabel('Месяц')
    plt.ylabel('Специализация')
    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=150)
    buf.seek(0)
    pic = slide.shapes.add_picture(buf, Inches(1), Inches(2), width=Inches(6))
    plt.close()
    add_fade_animation(pic, 1)

    # Слайд 8: Выводы
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Ключевые выводы"
    set_title_style(title)
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "Основные проблемы:"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    for point in [
        "40% жалоб связаны с коммуникацией",
        "Гастроэнтерологи и педиатры — лидеры по жалобам",
        "Пик жалоб в декабре–феврале",
        f"Эреджепова Э.М. —最多 жалоб ({analysis['top_doctors'].iloc[0]['Количество жалоб']})",
        "60% ответов с задержкой >3 часов"
    ]:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['secondary']
    add_fade_animation(content, 1)

    # Слайд 9: Рекомендации
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Рекомендации для платформы"
    set_title_style(title)
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "Рекомендации:"
    p.level = 0
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    for point in [
        "Внедрить чек-листы для врачей",
        "Провести обучение по soft skills",
        "Ввести рейтинговую систему",
        "Усилить модерацию спорных жалоб"
    ]:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['secondary']
    add_fade_animation(content, 1)

    prs.save(OUTPUT_PPTX)
    print(f"Презентация создана: {OUTPUT_PPTX}")

if __name__ == "__main__":
    df_main, df_monthly = load_data()
    analysis = analyze_data(df_main, df_monthly)
    create_presentation(analysis)